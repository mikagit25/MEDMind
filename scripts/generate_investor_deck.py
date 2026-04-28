"""
MedMind AI — Investor Deck Generator
Generates a professional PowerPoint (.pptx) investor presentation.

Usage:
    pip install python-pptx pillow
    python scripts/generate_investor_deck.py
    # Output: MedMind_Investor_Deck_2026.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

# ── Brand colours ──────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x1a, 0x18, 0x14)   # dark background
C_SURFACE  = RGBColor(0x2a, 0x25, 0x20)   # card surface
C_INK      = RGBColor(0xf0, 0xed, 0xe8)   # primary text (light)
C_INK2     = RGBColor(0xb8, 0xb0, 0xa4)   # secondary text
C_INK3     = RGBColor(0x7a, 0x72, 0x68)   # muted text
C_RED      = RGBColor(0xc0, 0x39, 0x2b)   # primary accent
C_RED2     = RGBColor(0xe7, 0x4c, 0x3c)   # lighter red
C_GREEN    = RGBColor(0x27, 0xae, 0x60)   # green highlight
C_BLUE     = RGBColor(0x29, 0x80, 0xb9)   # blue highlight
C_AMBER    = RGBColor(0xd4, 0xa0, 0x30)   # amber highlight
C_BORDER   = RGBColor(0x3d, 0x37, 0x30)   # card border / subtle lines
C_WHITE    = RGBColor(0xff, 0xff, 0xff)


# ── Slide dimensions — 16:9 widescreen ────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    """Add a fully blank slide (no placeholders)."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    return slide


def fill_bg(slide, color: RGBColor = C_BG):
    """Fill slide background with solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width_pt=0):
    """Add a filled rectangle shape."""
    from pptx.util import Pt as PtU
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()  # no line by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    return shape


def add_text(slide, text: str, left, top, width, height,
             font_size=18, bold=False, color=C_INK,
             align=PP_ALIGN.LEFT, word_wrap=True, italic=False):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txBox


def add_multiline(slide, lines: list, left, top, width, height,
                  font_size=14, color=C_INK2, line_spacing_pt=6, bold_first=False):
    """Add a text box with multiple bullet lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after = Pt(line_spacing_pt)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
        run.font.name = "Arial"
    return txBox


def slide_label(slide, text: str):
    """Small uppercase section label top-right."""
    add_text(slide, text,
             left=Inches(10.5), top=Inches(0.2), width=Inches(2.6), height=Inches(0.35),
             font_size=9, color=C_INK3, align=PP_ALIGN.RIGHT)


def logo(slide):
    """MedMind logo top-left."""
    txBox = slide.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(3), Inches(0.45))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "Med"
    run1.font.size = Pt(18)
    run1.font.bold = True
    run1.font.color.rgb = C_INK
    run1.font.name = "Arial"
    run2 = p.add_run()
    run2.text = "Mind"
    run2.font.size = Pt(18)
    run2.font.bold = True
    run2.font.color.rgb = C_RED2
    run2.font.name = "Arial"
    run3 = p.add_run()
    run3.text = " AI"
    run3.font.size = Pt(11)
    run3.font.bold = False
    run3.font.color.rgb = C_INK3
    run3.font.name = "Arial"


def divider(slide, y=Inches(0.72)):
    """Thin horizontal rule across the slide."""
    add_rect(slide, Inches(0.45), y, Inches(12.4), Pt(1),
             fill_color=C_BORDER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_cover(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    # Red accent bar left
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), H, fill_color=C_RED)

    # Logo large
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "Med"; r1.font.size = Pt(60); r1.font.bold = True
    r1.font.color.rgb = C_INK; r1.font.name = "Arial"
    r2 = p.add_run(); r2.text = "Mind"; r2.font.size = Pt(60); r2.font.bold = True
    r2.font.color.rgb = C_RED2; r2.font.name = "Arial"
    r3 = p.add_run(); r3.text = " AI"; r3.font.size = Pt(30); r3.font.bold = False
    r3.font.color.rgb = C_INK3; r3.font.name = "Arial"

    add_text(slide, "The Global AI Platform for Medical Education",
             Inches(0.8), Inches(3.0), Inches(9), Inches(0.7),
             font_size=22, color=C_INK2)

    add_text(slide, "Investor Presentation · 2026",
             Inches(0.8), Inches(3.85), Inches(6), Inches(0.45),
             font_size=13, color=C_INK3)

    # Bottom tag line
    add_rect(slide, Inches(0.8), Inches(6.3), Inches(5.5), Inches(0.6),
             fill_color=C_SURFACE)
    add_text(slide, "Confidential — Not for distribution",
             Inches(0.85), Inches(6.35), Inches(5.4), Inches(0.5),
             font_size=11, color=C_INK3)

    # Right side — key numbers
    for i, (val, lbl) in enumerate([
        ("$6.4B", "Market size (2024)"),
        ("18%", "Annual market growth"),
        ("12M+", "Target clinicians"),
        ("82+", "Modules built today"),
    ]):
        y = Inches(1.8 + i * 1.2)
        add_rect(slide, Inches(9.8), y, Inches(3.0), Inches(0.95),
                 fill_color=C_SURFACE)
        add_text(slide, val,  Inches(9.9), y + Inches(0.05), Inches(2.8), Inches(0.45),
                 font_size=26, bold=True, color=C_RED2)
        add_text(slide, lbl,  Inches(9.9), y + Inches(0.48), Inches(2.8), Inches(0.38),
                 font_size=11, color=C_INK3)


def slide_problem(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    logo(slide); divider(slide); slide_label(slide, "01 / Problem")

    add_text(slide, "Medical education is broken",
             Inches(0.45), Inches(0.9), Inches(12), Inches(0.7),
             font_size=30, bold=True, color=C_INK)

    add_text(slide, "Clinicians need to learn continuously — but the tools are fragmented, expensive, and outdated.",
             Inches(0.45), Inches(1.65), Inches(12), Inches(0.5),
             font_size=14, color=C_INK2)

    problems = [
        ("📖", "Static textbooks", "Outdated the moment they're printed. No adaptation, no interactivity."),
        ("🔀", "Fragmented tools", "Flashcard app + chatbot + PDF textbook + YouTube. No unified system."),
        ("💸", "Prohibitive cost", "Quality board prep courses cost $500–$2,000. Inaccessible to most."),
        ("🌐", "English-only", "Best tools exclude 70%+ of global medical students."),
        ("🤖", "Generic AI", "ChatGPT lacks clinical structure, cites nothing, can't replace a curriculum."),
    ]

    for i, (icon, title, desc) in enumerate(problems):
        col = i % 3
        row = i // 3
        x = Inches(0.45 + col * 4.3)
        y = Inches(2.4 + row * 2.0)
        add_rect(slide, x, y, Inches(4.0), Inches(1.7), fill_color=C_SURFACE)
        add_text(slide, icon,  x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.45), font_size=20)
        add_text(slide, title, x + Inches(0.7), y + Inches(0.15), Inches(3.1), Inches(0.4),
                 font_size=13, bold=True, color=C_INK)
        add_text(slide, desc,  x + Inches(0.2), y + Inches(0.65), Inches(3.6), Inches(0.9),
                 font_size=11, color=C_INK3)


def slide_solution(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    logo(slide); divider(slide); slide_label(slide, "02 / Solution")

    add_text(slide, "One complete AI-native platform",
             Inches(0.45), Inches(0.9), Inches(12), Inches(0.7),
             font_size=30, bold=True, color=C_INK)

    solutions = [
        ("🧠", "AI Medical Tutor",     "4 modes: Tutor · Socratic · Case-based · Exam Prep.\nEvery answer backed by live PubMed citations."),
        ("📇", "Spaced Repetition",    "SM-2 algorithm — the same science as Anki.\n500+ medical flashcards, auto-scheduled."),
        ("🩺", "Clinical Simulations", "Interactive patient cases: vitals, labs, differentials.\nUSMLE / MRCP exam formats."),
        ("📚", "82+ Modules",          "Cardiology · Neurology · Surgery · Pediatrics · OB/GYN · Therapy.\nFull structured curriculum."),
        ("🌍", "7 Languages",          "English · Russian · German · French · Spanish · Turkish · Arabic.\nBuilt-in from day one."),
        ("📈", "Progress Analytics",   "Streak calendar, retention scores, completion rates.\nFor students and institutional admins."),
    ]

    for i, (icon, title, desc) in enumerate(solutions):
        col = i % 3
        row = i // 3
        x = Inches(0.45 + col * 4.3)
        y = Inches(1.7 + row * 2.45)
        add_rect(slide, x, y, Inches(4.05), Inches(2.2), fill_color=C_SURFACE)
        add_text(slide, icon,  x + Inches(0.2), y + Inches(0.18), Inches(0.5), Inches(0.45), font_size=22)
        add_text(slide, title, x + Inches(0.75), y + Inches(0.18), Inches(3.1), Inches(0.4),
                 font_size=13, bold=True, color=C_INK)
        add_text(slide, desc,  x + Inches(0.2), y + Inches(0.72), Inches(3.65), Inches(1.3),
                 font_size=10.5, color=C_INK3)


def slide_market(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    logo(slide); divider(slide); slide_label(slide, "03 / Market")

    add_text(slide, "A large, growing, underserved market",
             Inches(0.45), Inches(0.9), Inches(12), Inches(0.7),
             font_size=30, bold=True, color=C_INK)

    # Three circles / TAM-SAM-SOM
    circles = [
        ("TAM", "$6.4B", "Global medical\ne-learning market"),
        ("SAM", "$1.2B", "Digital-first students\n& residents globally"),
        ("SOM", "$120M", "Year 5 addressable\nwith current product"),
    ]
    colors = [C_RED, RGBColor(0xc0, 0x60, 0x50), RGBColor(0xb0, 0x80, 0x70)]
    sizes  = [Inches(3.8), Inches(3.2), Inches(2.6)]
    for i, ((label, val, desc), color, sz) in enumerate(zip(circles, colors, sizes)):
        cx = Inches(2.0 + i * 3.8)
        cy = Inches(2.5)
        # Ellipse
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        shape = slide.shapes.add_shape(9, cx, cy, sz, sz)  # 9 = oval
        shape.fill.solid(); shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        # Label
        add_text(slide, label, cx, cy + sz * 0.1, sz, Inches(0.4),
                 font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, val,   cx, cy + sz * 0.35, sz, Inches(0.55),
                 font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc,  cx, cy + sz * 0.62, sz, Inches(0.6),
                 font_size=10, color=RGBColor(0xff, 0xdd, 0xd8), align=PP_ALIGN.CENTER)

    # Stats row
    stats = [
        ("~2M",   "Medical students\nglobally"),
        ("~10M",  "Practicing physicians"),
        ("$299–$699", "Annual study tool spend\nper student"),
        ("18%",   "Market CAGR\n(2024–2029)"),
    ]
    for i, (val, lbl) in enumerate(stats):
        x = Inches(0.45 + i * 3.2)
        y = Inches(5.9)
        add_rect(slide, x, y, Inches(3.0), Inches(1.15), fill_color=C_SURFACE)
        add_text(slide, val, x + Inches(0.15), y + Inches(0.08), Inches(2.7), Inches(0.5),
                 font_size=22, bold=True, color=C_RED2)
        add_text(slide, lbl, x + Inches(0.15), y + Inches(0.6), Inches(2.7), Inches(0.45),
                 font_size=10, color=C_INK3)


def slide_traction(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    logo(slide); divider(slide); slide_label(slide, "04 / Traction")

    add_text(slide, "Platform built. Ready to deploy.",
             Inches(0.45), Inches(0.9), Inches(12), Inches(0.7),
             font_size=30, bold=True, color=C_INK)

    add_text(slide, "Not a pitch deck — a working product.",
             Inches(0.45), Inches(1.6), Inches(8), Inches(0.4),
             font_size=14, color=C_INK3)

    done = [
        "82+ clinical modules across 6 specialties",
        "Full AI tutor — 4 learning modes",
        "Spaced-repetition flashcard engine (SM-2)",
        "Clinical case simulation system",
        "7-language localisation (en/ru/de/fr/es/tr/ar)",
        "SEO article generation engine (Claude AI)",
        "Article-to-video pipeline (YouTube-ready)",
        "B2B Clinic tier with team analytics",
        "Admin panel & content management",
        "Full authentication + role system",
    ]
    next_milestones = [
        "Production server deployment",
        "SEO content campaign (500+ articles)",
        "YouTube medical channel launch",
        "First 100 paying subscribers",
        "University pilot program",
        "Series A fundraise",
    ]

    add_rect(slide, Inches(0.45), Inches(2.1), Inches(6.0), Inches(4.9), fill_color=C_SURFACE)
    add_text(slide, "✅  Already done", Inches(0.6), Inches(2.2), Inches(5.7), Inches(0.4),
             font_size=13, bold=True, color=C_GREEN)
    for i, item in enumerate(done):
        add_text(slide, f"✓  {item}", Inches(0.6), Inches(2.65 + i * 0.42),
                 Inches(5.7), Inches(0.38), font_size=11, color=C_INK2)

    add_rect(slide, Inches(6.9), Inches(2.1), Inches(6.0), Inches(4.9), fill_color=C_SURFACE)
    add_text(slide, "◎  Next milestones", Inches(7.05), Inches(2.2), Inches(5.7), Inches(0.4),
             font_size=13, bold=True, color=C_AMBER)
    for i, item in enumerate(next_milestones):
        add_text(slide, f"→  {item}", Inches(7.05), Inches(2.65 + i * 0.42),
                 Inches(5.7), Inches(0.38), font_size=11, color=C_INK2)


def slide_business_model(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    logo(slide); divider(slide); slide_label(slide, "05 / Business Model")

    add_text(slide, "Three revenue streams",
             Inches(0.45), Inches(0.9), Inches(12), Inches(0.7),
             font_size=30, bold=True, color=C_INK)

    streams = [
        {
            "title": "B2C Subscriptions",
            "icon": "💳",
            "color": C_RED,
            "tiers": [
                ("Free",     "$0",       "8 modules · 5 AI Q/day"),
                ("Student",  "$15/mo",   "82+ modules · 50 AI Q/day"),
                ("Pro",      "$40/mo",   "Unlimited AI · Drug DB"),
                ("Lifetime", "$299 once","All future content"),
            ],
            "note": "Primary engine. High LTV, low CAC via SEO.",
        },
        {
            "title": "B2B Institutional",
            "icon": "🏫",
            "color": C_BLUE,
            "tiers": [
                ("Clinic",     "$199/mo", "Up to 10 seats + analytics"),
                ("Enterprise", "Custom",  "Medical schools & hospitals"),
            ],
            "note": "Sticky annual contracts. 1 school = 200–500 seats.",
        },
        {
            "title": "Content & Licensing",
            "icon": "📄",
            "color": C_AMBER,
            "tiers": [
                ("SEO Articles", "Ad/affiliate", "500+ indexed medical articles"),
                ("YouTube",      "Ad revenue",   "AI-generated multilingual videos"),
                ("API",          "Usage-based",  "Medical AI Q&A for EHR/health-tech"),
            ],
            "note": "Long-term diversification beyond subscriptions.",
        },
    ]

    for i, s in enumerate(streams):
        x = Inches(0.45 + i * 4.3)
        add_rect(slide, x, Inches(1.7), Inches(4.05), Inches(5.45), fill_color=C_SURFACE)
        # colour top bar
        add_rect(slide, x, Inches(1.7), Inches(4.05), Inches(0.07), fill_color=s["color"])
        add_text(slide, s["icon"] + "  " + s["title"],
                 x + Inches(0.2), Inches(1.8), Inches(3.7), Inches(0.45),
                 font_size=14, bold=True, color=C_INK)
        for j, (tier, price, desc) in enumerate(s["tiers"]):
            y = Inches(2.4 + j * 1.05)
            add_rect(slide, x + Inches(0.15), y, Inches(3.7), Inches(0.9),
                     fill_color=C_BG)
            add_text(slide, tier,  x + Inches(0.3), y + Inches(0.05), Inches(2.0), Inches(0.35),
                     font_size=12, bold=True, color=C_INK)
            add_text(slide, price, x + Inches(0.3), y + Inches(0.45), Inches(1.5), Inches(0.35),
                     font_size=13, bold=True, color=s["color"])
            add_text(slide, desc,  x + Inches(1.9), y + Inches(0.45), Inches(1.9), Inches(0.35),
                     font_size=9.5, color=C_INK3)
        add_text(slide, s["note"], x + Inches(0.2), Inches(6.55), Inches(3.7), Inches(0.5),
                 font_size=10, color=C_INK3, italic=True)


def slide_unit_economics(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    logo(slide); divider(slide); slide_label(slide, "06 / Unit Economics")

    add_text(slide, "Strong LTV / CAC dynamics",
             Inches(0.45), Inches(0.9), Inches(12), Inches(0.7),
             font_size=30, bold=True, color=C_INK)

    metrics = [
        ("Student LTV",        "$360",   "24 mo avg retention @ $15/mo",    C_GREEN),
        ("Pro LTV",            "$720",   "18 mo avg retention @ $40/mo",     C_GREEN),
        ("Clinic LTV",         "$2,388", "12 mo avg contract @ $199/mo",     C_GREEN),
        ("Lifetime order",     "$299",   "One-time, full margin",             C_GREEN),
        ("CAC — SEO",          "< $10",  "Organic-first content strategy",    C_BLUE),
        ("CAC — Paid social",  "< $30",  "Target blended CAC",               C_BLUE),
        ("Gross margin",       "~85%",   "SaaS software margin",              C_AMBER),
        ("Payback period",     "< 1 mo", "Student plan @ $10 CAC",           C_AMBER),
    ]

    for i, (label, val, note, color) in enumerate(metrics):
        col = i % 4
        row = i // 4
        x = Inches(0.45 + col * 3.22)
        y = Inches(1.85 + row * 2.3)
        add_rect(slide, x, y, Inches(3.0), Inches(2.0), fill_color=C_SURFACE)
        add_rect(slide, x, y, Inches(3.0), Inches(0.06), fill_color=color)
        add_text(slide, val,   x + Inches(0.2), y + Inches(0.15), Inches(2.6), Inches(0.6),
                 font_size=28, bold=True, color=color)
        add_text(slide, label, x + Inches(0.2), y + Inches(0.78), Inches(2.6), Inches(0.45),
                 font_size=12, bold=True, color=C_INK)
        add_text(slide, note,  x + Inches(0.2), y + Inches(1.25), Inches(2.6), Inches(0.6),
                 font_size=10, color=C_INK3)


def slide_moat(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    logo(slide); divider(slide); slide_label(slide, "07 / Competitive Moat")

    add_text(slide, "Why MedMind wins",
             Inches(0.45), Inches(0.9), Inches(12), Inches(0.7),
             font_size=30, bold=True, color=C_INK)

    moats = [
        ("🔬", "PubMed-backed AI",       "Real-time evidence in every answer.\nStructural advantage vs. generic AI."),
        ("📚", "Curriculum depth",        "82+ structured modules — years of work.\nNot replicable by copy-paste."),
        ("🌍", "Multilingual at launch",  "7 languages from day one.\nCompetitors are English-only."),
        ("📈", "SEO content flywheel",    "Hundreds of medical articles drive organic traffic.\nCompounding without paid ads."),
        ("🔄", "Daily habit loop",        "Spaced repetition + AI = users return daily.\nBuilt-in retention mechanism."),
        ("👥", "Multi-role platform",     "Student · Resident · Doctor · Professor · Vet.\nEnables institutional contracts."),
    ]

    for i, (icon, title, desc) in enumerate(moats):
        col = i % 3
        row = i // 3
        x = Inches(0.45 + col * 4.3)
        y = Inches(1.75 + row * 2.45)
        add_rect(slide, x, y, Inches(4.05), Inches(2.25), fill_color=C_SURFACE)
        add_text(slide, icon,  x + Inches(0.2), y + Inches(0.18), Inches(0.5), Inches(0.45), font_size=22)
        add_text(slide, title, x + Inches(0.75), y + Inches(0.18), Inches(3.1), Inches(0.4),
                 font_size=13, bold=True, color=C_INK)
        add_text(slide, desc,  x + Inches(0.2), y + Inches(0.75), Inches(3.65), Inches(1.3),
                 font_size=11, color=C_INK3)


def slide_roadmap(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    logo(slide); divider(slide); slide_label(slide, "08 / Roadmap")

    add_text(slide, "Three phases to scale",
             Inches(0.45), Inches(0.9), Inches(12), Inches(0.7),
             font_size=30, bold=True, color=C_INK)

    phases = [
        {
            "phase": "Phase 1",
            "period": "Q2–Q3 2026",
            "title": "Launch & Revenue",
            "color": C_RED,
            "items": [
                "Production server deployment",
                "500+ SEO medical articles live",
                "YouTube channel with medical video content",
                "App Store / Google Play listing",
                "100 paying subscribers",
            ],
        },
        {
            "phase": "Phase 2",
            "period": "Q4 2026–Q1 2027",
            "title": "B2B & Institutions",
            "color": C_BLUE,
            "items": [
                "First university pilot (100–300 seats)",
                "SCORM / LTI LMS integration",
                "White-label offering for medical schools",
                "CME / CPD credit integration",
                "1,000+ monthly active users",
            ],
        },
        {
            "phase": "Phase 3",
            "period": "2027+",
            "title": "Scale & API",
            "color": C_GREEN,
            "items": [
                "Medical AI API for EHR / health-tech",
                "15+ languages",
                "National medical board partnerships",
                "Series A fundraise",
                "10,000 paying users B2C + B2B",
            ],
        },
    ]

    # Timeline line
    add_rect(slide, Inches(0.45), Inches(2.15), Inches(12.4), Pt(2), fill_color=C_BORDER)

    for i, ph in enumerate(phases):
        x = Inches(0.45 + i * 4.3)
        # Dot on timeline
        dot = slide.shapes.add_shape(9, x + Inches(1.5), Inches(2.0), Inches(0.3), Inches(0.3))
        dot.fill.solid(); dot.fill.fore_color.rgb = ph["color"]
        dot.line.fill.background()

        add_rect(slide, x, Inches(2.55), Inches(4.05), Inches(4.6), fill_color=C_SURFACE)
        add_rect(slide, x, Inches(2.55), Inches(4.05), Inches(0.07), fill_color=ph["color"])

        add_text(slide, ph["phase"] + "  ·  " + ph["period"],
                 x + Inches(0.2), Inches(2.65), Inches(3.7), Inches(0.35),
                 font_size=10, bold=True, color=ph["color"])
        add_text(slide, ph["title"],
                 x + Inches(0.2), Inches(3.05), Inches(3.7), Inches(0.45),
                 font_size=15, bold=True, color=C_INK)
        for j, item in enumerate(ph["items"]):
            add_text(slide, "→  " + item,
                     x + Inches(0.2), Inches(3.6 + j * 0.66), Inches(3.7), Inches(0.55),
                     font_size=11, color=C_INK3)


def slide_ask(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    logo(slide); divider(slide); slide_label(slide, "09 / The Ask")

    add_text(slide, "Pre-seed round",
             Inches(0.45), Inches(0.9), Inches(12), Inches(0.7),
             font_size=30, bold=True, color=C_INK)

    # Main ask box
    add_rect(slide, Inches(0.45), Inches(1.75), Inches(5.8), Inches(5.3), fill_color=C_SURFACE)
    add_rect(slide, Inches(0.45), Inches(1.75), Inches(5.8), Inches(0.07), fill_color=C_RED)

    add_text(slide, "Raising", Inches(0.65), Inches(1.9), Inches(5.4), Inches(0.4),
             font_size=12, color=C_INK3, bold=True)
    add_text(slide, "$150K – $300K",
             Inches(0.65), Inches(2.35), Inches(5.4), Inches(0.7),
             font_size=32, bold=True, color=C_RED2)
    add_text(slide, "Pre-seed · SAFE / Convertible note",
             Inches(0.65), Inches(3.1), Inches(5.4), Inches(0.4),
             font_size=12, color=C_INK3)

    add_text(slide, "Use of funds", Inches(0.65), Inches(3.65), Inches(5.4), Inches(0.4),
             font_size=12, bold=True, color=C_INK)
    for j, (pct, item) in enumerate([
        ("40%", "Infrastructure & DevOps — production servers, CDN, monitoring"),
        ("30%", "Marketing launch — SEO content, paid acquisition, App Store"),
        ("20%", "Product development — mobile app, API, new specialties"),
        ("10%", "Legal & operations — company formation, contracts"),
    ]):
        y = Inches(4.1 + j * 0.55)
        add_rect(slide, Inches(0.65), y + Inches(0.08), Inches(0.55), Inches(0.32),
                 fill_color=C_RED)
        add_text(slide, pct, Inches(0.65), y + Inches(0.05), Inches(0.55), Inches(0.38),
                 font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, item, Inches(1.28), y, Inches(4.6), Inches(0.5),
                 font_size=10.5, color=C_INK2)

    # Contact box
    add_rect(slide, Inches(6.8), Inches(1.75), Inches(6.0), Inches(5.3), fill_color=C_SURFACE)
    add_rect(slide, Inches(6.8), Inches(1.75), Inches(6.0), Inches(0.07), fill_color=C_GREEN)

    add_text(slide, "Get in touch", Inches(7.0), Inches(1.9), Inches(5.6), Inches(0.45),
             font_size=16, bold=True, color=C_INK)

    contacts = [
        ("Investment enquiries",   "invest@medmind.pro"),
        ("Partnership / B2B",       "partners@medmind.pro"),
        ("Website",                 "medmind.pro"),
        ("Try the product",         "medmind.pro/register"),
    ]
    for j, (label, val) in enumerate(contacts):
        y = Inches(2.5 + j * 0.9)
        add_rect(slide, Inches(7.0), y, Inches(5.6), Inches(0.75), fill_color=C_BG)
        add_text(slide, label, Inches(7.15), y + Inches(0.05), Inches(5.3), Inches(0.3),
                 font_size=10, color=C_INK3)
        add_text(slide, val,   Inches(7.15), y + Inches(0.35), Inches(5.3), Inches(0.35),
                 font_size=13, bold=True, color=C_RED2)

    add_text(slide, "The platform is live. All we need is fuel.",
             Inches(7.0), Inches(6.1), Inches(5.6), Inches(0.5),
             font_size=13, color=C_INK3, italic=True)


def slide_closing(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    add_rect(slide, Inches(0), Inches(0), Inches(0.18), H, fill_color=C_RED)

    add_text(slide, "Thank you",
             Inches(0.8), Inches(1.8), Inches(11), Inches(1.1),
             font_size=52, bold=True, color=C_INK, align=PP_ALIGN.CENTER)

    add_text(slide, "MedMind AI — Evidence-based learning for every clinician, everywhere.",
             Inches(0.8), Inches(3.0), Inches(11), Inches(0.7),
             font_size=18, color=C_INK2, align=PP_ALIGN.CENTER)

    add_text(slide, "invest@medmind.pro  ·  medmind.pro",
             Inches(0.8), Inches(3.9), Inches(11), Inches(0.5),
             font_size=14, color=C_RED2, align=PP_ALIGN.CENTER)

    add_text(slide, "Confidential — For discussion purposes only. Not a securities offering.",
             Inches(1.5), Inches(6.6), Inches(10), Inches(0.4),
             font_size=9, color=C_INK3, align=PP_ALIGN.CENTER)

    # Logo large centred
    txBox = slide.shapes.add_textbox(Inches(4.5), Inches(4.7), Inches(4), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = "Med"; r1.font.size = Pt(32); r1.font.bold = True
    r1.font.color.rgb = C_INK; r1.font.name = "Arial"
    r2 = p.add_run(); r2.text = "Mind"; r2.font.size = Pt(32); r2.font.bold = True
    r2.font.color.rgb = C_RED2; r2.font.name = "Arial"
    r3 = p.add_run(); r3.text = " AI"; r3.font.size = Pt(18)
    r3.font.color.rgb = C_INK3; r3.font.name = "Arial"


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    print("Building slides…")
    slide_cover(prs)
    print("  01 Cover")
    slide_problem(prs)
    print("  02 Problem")
    slide_solution(prs)
    print("  03 Solution")
    slide_market(prs)
    print("  04 Market")
    slide_traction(prs)
    print("  05 Traction")
    slide_business_model(prs)
    print("  06 Business Model")
    slide_unit_economics(prs)
    print("  07 Unit Economics")
    slide_moat(prs)
    print("  08 Competitive Moat")
    slide_roadmap(prs)
    print("  09 Roadmap")
    slide_ask(prs)
    print("  10 The Ask")
    slide_closing(prs)
    print("  11 Closing")

    out = "MedMind_Investor_Deck_2026.pptx"
    prs.save(out)
    print(f"\nSaved → {out}")
    print("Open in PowerPoint or Google Slides.")
    print("To export PDF: File → Export → PDF (PowerPoint) or print-to-PDF.")


if __name__ == "__main__":
    main()
